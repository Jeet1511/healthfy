from pptx import Presentation

ppt_path = r"c:\Users\Jeet\Desktop\hackathon\OMINA_Hackathon_12_Slides_V2.pptx"
prs = Presentation(ppt_path)

for i, slide in enumerate(prs.slides, start=1):
    print(f"\n--- Slide {i} ---")
    for j, shape in enumerate(slide.shapes, start=1):
        has_text = hasattr(shape, "text")
        text = shape.text.strip().replace("\n", " | ") if has_text and shape.text else ""
        shape_type = getattr(shape, "shape_type", None)
        placeholder_type = None
        if getattr(shape, "is_placeholder", False):
            try:
                placeholder_type = str(shape.placeholder_format.type)
            except Exception:
                placeholder_type = "unknown"
        print(f"{j:02d}. placeholder={shape.is_placeholder} ptype={placeholder_type} type={shape_type} text='{text[:120]}'")
