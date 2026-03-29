from pptx import Presentation

ppt_path = r"c:\Users\Jeet\Desktop\hackathon\OMINA_Hackathon_12_Slides_V2.pptx"
prs = Presentation(ppt_path)

slides_content = [
    {
        "title": "TITLE PAGE",
        "body": [
            "Team Name: Zenith-X",
            "Project Name: OMINA — Universal Human Assistance Platform",
            "Team Leader: Jeet Mondal",
            "Members: [Add all member names]",
            "Institute Name: [Add institute name]",
            "Hackathon: [Add hackathon name/year]",
        ],
    },
    {
        "title": "IDEA TITLE",
        "body": [
            "Proposed Solution (Describe your Idea/Solution/Prototype)",
            "Detailed explanation of the proposed solution: OMINA is a dual-mode assistance platform for daily needs and emergency response in one app.",
            "How it addresses the problem: It removes fragmentation and reduces delay by giving context-aware, priority-based actions.",
            "Innovation and uniqueness of the solution: Context-adaptive UX (Daily ↔ Emergency) with AI-backed urgency classification.",
        ],
    },
    {
        "title": "TECHNICAL APPROACH",
        "body": [
            "Technologies to be used (e.g. programming languages, frameworks, hardware): MERN stack, Vite, JWT Auth, Axios, AI provider gateway.",
            "Methodology and process for implementation (Flow Charts/Images/ working prototype): Requirement mapping → dual-mode UI design → modular API build → intent/risk integration → testing → working prototype.",
        ],
    },
    {
        "title": "TECHNICAL APPROACH",
        "body": [
            "For flowcharts and Diagrams",
            "System architecture diagram, intent-to-action flowchart, emergency escalation flow, and deployment diagram.",
        ],
    },
    {
        "title": "FEASIBILITY AND VIABILITY",
        "body": [
            "Analysis of the feasibility of the idea: Technically feasible with existing MERN architecture and modular services.",
            "Potential challenges and risks: AI latency, data privacy/security, and local provider dependency.",
            "Strategies for overcoming these challenges: Fallback logic, secure auth/encryption, and phased city-wise provider onboarding.",
        ],
    },
    {
        "title": "IMPACT AND BENEFITS",
        "body": [
            "Potential impact on the target audience: Faster, clearer access to help for students, women, commuters, and families.",
            "Benefits of the solution (social, economic, environmental, etc.): Improved safety outcomes, reduced response inefficiency, and stronger community resilience.",
        ],
    },
    {
        "title": "MARKET ANALYSIS",
        "body": [
            "Define the specific user niche (Beachhead) and the total potential market size (TAM) to prove scalability: Beachhead is urban students and working women; TAM spans smart-city digital safety and assistance users.",
            "State your unique technical or operational advantage that prevents competitors from easily replicating your solution: Dual-UX architecture with AI decision normalization and multi-provider fallback.",
            "State the business model: B2B2C subscriptions with institutions plus premium consumer safety features and partner analytics APIs.",
        ],
    },
]


def get_title_placeholder(slide):
    for shape in slide.shapes:
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type in (1, 3):  # TITLE or CENTER_TITLE
                    return shape
            except Exception:
                pass
    return None


def get_body_placeholder(slide):
    for shape in slide.shapes:
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type in (7, 4):  # OBJECT/BODY or SUBTITLE
                    return shape
            except Exception:
                pass
    return None


# Ensure exactly 7 slides
sld_id_lst = prs.slides._sldIdLst
while len(prs.slides) > len(slides_content):
    sld_id_lst.remove(sld_id_lst[-1])

for idx, content in enumerate(slides_content):
    slide = prs.slides[idx]

    title_shape = get_title_placeholder(slide)
    if title_shape is not None:
        title_shape.text = content["title"]

    body_shape = get_body_placeholder(slide)
    if body_shape is not None:
        tf = body_shape.text_frame
        tf.clear()

        for i, line in enumerate(content["body"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.level = 0

prs.save(ppt_path)
print(f"Updated PPT: {ppt_path}")
print(f"Total slides: {len(prs.slides)}")
