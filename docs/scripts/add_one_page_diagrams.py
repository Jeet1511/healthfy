from pptx import Presentation
from pathlib import Path

ppt_path = Path(r"c:\Users\Jeet\Desktop\hackathon\OMINA_Hackathon_12_Slides_V2.pptx")
prs = Presentation(str(ppt_path))

# Slide 4 in user view is index 3 (0-based)
slide = prs.slides[3]

# Find title + body placeholders
title_shape = None
body_shape = None
for shape in slide.shapes:
    if shape.is_placeholder:
        ptype = shape.placeholder_format.type
        if ptype in (1, 3):  # TITLE / CENTER_TITLE
            title_shape = shape
        elif ptype in (7, 4):  # OBJECT / SUBTITLE
            body_shape = shape

if title_shape is not None:
    title_shape.text = "TECHNICAL APPROACH"

diagram_text = (
    "ONE-PAGE FLOWCHARTS & DIAGRAMS\n"
    "\n"
    "1) SYSTEM ARCHITECTURE\n"
    "User App (React) -> API Gateway (Express) -> Decision Engine\n"
    "Decision Engine -> AI Providers (OpenAI | Anthropic | Groq)\n"
    "API Gateway <-> MongoDB (Users, Requests, Providers)\n"
    "\n"
    "2) INTENT-TO-ACTION FLOW\n"
    "User Input -> Intent Classification -> Risk Level (CRITICAL/MODERATE/SAFE)\n"
    "-> Mode Switch (Emergency/Daily) -> Action Cards + Instructions\n"
    "\n"
    "3) EMERGENCY ESCALATION FLOW\n"
    "Distress Query -> CRITICAL Tag -> SOS + Safe Zones + Live Tracking\n"
    "-> Assistance Request Logged -> Follow-up Notifications\n"
    "\n"
    "4) DEPLOYMENT DIAGRAM\n"
    "Frontend: Vercel  <->  Backend API: Cloud Node Service  <->  MongoDB Atlas\n"
    "(Secure JWT Auth + Environment Config + Provider Key Management)"
)

if body_shape is not None:
    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True

    lines = diagram_text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0

try:
    prs.save(str(ppt_path))
    print(f"Updated one-page diagrams on slide 4: {ppt_path}")
except PermissionError:
    output_path = ppt_path.with_name("OMINA_Hackathon_12_Slides_V3.pptx")
    prs.save(str(output_path))
    print(f"Original file locked. Saved updated deck as: {output_path}")
